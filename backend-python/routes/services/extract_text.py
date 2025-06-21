from pdfminer.high_level import extract_text as pdf_extract_text
from pptx import Presentation

def extract_text_from_pdf(file_path: str) -> str:
    return pdf_extract_text(file_path)


def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                texts.append(shape.text)
    return '\n'.join(texts)
